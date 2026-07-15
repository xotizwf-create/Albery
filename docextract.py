"""Text extraction for the long tail of document formats the agents receive.

b24bot._b24_extract_document keeps the primary formats (pdf/docx/xlsx/plain text) and
delegates everything else here: legacy Office (.xls/.doc/.ppt), presentations (.pptx),
OpenDocument (.odt/.ods/.odp), web archives (.mht/.mhtml), HTML, RTF and ZIP containers.

Contract: extract(data, name, inner=...) returns readable text or "" when the format
genuinely can't be read — the caller then stores the raw file and answers honestly.
Extractors never raise. Legacy binary formats use the tiny catdoc/antiword CLI tools
(apt: catdoc, antiword) via a temp file; everything else is pure python."""

from __future__ import annotations

import io
import logging
import os
import re
import subprocess
import tempfile
import zipfile

MAX_CHARS = 600_000            # same budget as webread.extract_xlsx
_ZIP_MAX_MEMBERS = 8           # how many readable files inside an archive to extract
_ZIP_MAX_MEMBER_BYTES = 15 * 1024 * 1024
_ZIP_MAX_TOTAL_BYTES = 60 * 1024 * 1024  # decompressed budget — zip-bomb guard

EXTS = ("xls", "pptx", "ppt", "doc", "mht", "mhtml", "htm", "html",
        "odt", "ods", "odp", "rtf", "zip", "xml")


def extract(data: bytes, name: str, inner=None) -> str:
    """Route by extension. `inner` is the primary extractor (b24bot._b24_extract_document),
    used for files found inside ZIP archives — passed in to avoid a circular import."""
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    try:
        if ext == "xls":
            return _xls(data)
        if ext == "pptx":
            return _pptx(data)
        if ext == "ppt":
            return _cli_tool(["catppt", "-dutf-8"], data, ".ppt")
        if ext == "doc":
            return _doc(data)
        if ext in ("mht", "mhtml"):
            return _mhtml(data)
        if ext in ("htm", "html"):
            return html_to_text(data.decode("utf-8", "ignore"))
        if ext in ("odt", "ods", "odp"):
            return _odf(data)
        if ext == "rtf":
            return _rtf(data)
        if ext == "zip":
            return _zip(data, inner)
        if ext == "xml":
            return data.decode("utf-8", "ignore")
    except Exception as exc:  # noqa: BLE001
        logging.warning("docextract failed (%s): %s", name, repr(exc)[:200])
    return ""


def html_to_text(html: str) -> str:
    """Visible text of an HTML page (lxml), with <script>/<style> dropped and
    blank-line runs collapsed — raw markup wastes the model's context."""
    if not (html or "").strip():
        return ""
    try:
        import lxml.html
        doc = lxml.html.fromstring(html)
        for bad in doc.xpath("//script | //style | //noscript"):
            bad.getparent().remove(bad)
        text = doc.text_content()
    except Exception:  # noqa: BLE001
        text = re.sub(r"<[^>]+>", " ", html)
    lines = [ln.strip() for ln in text.splitlines()]
    out, blank = [], 0
    for ln in lines:
        blank = blank + 1 if not ln else 0
        if blank <= 1:
            out.append(ln)
    return "\n".join(out).strip()[:MAX_CHARS]


def _xls(data: bytes) -> str:
    """Legacy Excel via xlrd. Mirrors webread.extract_xlsx's '# Лист:' + ' | ' table style."""
    import xlrd
    book = xlrd.open_workbook(file_contents=data)
    lines: list[str] = []
    total = 0
    for sheet in book.sheets():
        lines.append(f"# Лист: {sheet.name}")
        for r in range(sheet.nrows):
            vals = []
            for c in range(sheet.ncols):
                cell = sheet.cell(r, c)
                v = cell.value
                if cell.ctype == xlrd.XL_CELL_DATE:
                    try:
                        v = xlrd.xldate_as_datetime(v, book.datemode).isoformat(sep=" ")
                    except Exception:  # noqa: BLE001
                        pass
                elif cell.ctype == xlrd.XL_CELL_NUMBER and float(v).is_integer():
                    v = int(v)
                vals.append("" if v is None else str(v).strip())
            if any(vals):
                line = " | ".join(vals).rstrip(" |")
                lines.append(line)
                total += len(line)
                if total > MAX_CHARS:
                    lines.append("…[обрезано по лимиту]")
                    return "\n".join(lines)
    return "\n".join(lines).strip()


def _pptx(data: bytes) -> str:
    """Slide texts + tables + speaker notes via python-pptx."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"# Слайд {i}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for p in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in p.runs) or (p.text or "")
                    if t.strip():
                        lines.append(t.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    lines.append(" | ".join(cell.text.strip() for cell in row.cells))
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                lines.append("[Заметки докладчика] " + notes)
        if sum(len(x) for x in lines) > MAX_CHARS:
            lines.append("…[обрезано по лимиту]")
            break
    return "\n".join(lines).strip()


def _doc(data: bytes) -> str:
    """Legacy binary .doc: antiword first (better tables), catdoc as fallback."""
    return (_cli_tool(["antiword", "-m", "UTF-8.txt"], data, ".doc")
            or _cli_tool(["catdoc", "-dutf-8"], data, ".doc"))


def _cli_tool(cmd: list[str], data: bytes, suffix: str) -> str:
    """Run a tiny converter CLI (catdoc/catppt/antiword) over a temp copy of the file.
    Missing binary or failure -> "" (the caller falls back or reports honestly)."""
    path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tf:
            tf.write(data)
            path = tf.name
        r = subprocess.run(cmd + [path], capture_output=True, timeout=30)
        if r.returncode != 0:
            return ""
        return r.stdout.decode("utf-8", "ignore").strip()[:MAX_CHARS]
    except FileNotFoundError:
        logging.warning("docextract: converter not installed: %s", cmd[0])
        return ""
    except Exception as exc:  # noqa: BLE001
        logging.warning("docextract cli failed (%s): %s", cmd[0], repr(exc)[:160])
        return ""
    finally:
        if path:
            try:
                os.unlink(path)
            except OSError:
                pass


def _mhtml(data: bytes) -> str:
    """Saved web page (.mht/.mhtml): MIME multipart, main part is quoted-printable HTML.
    Browser snapshots (Chrome/Yandex 'Saved by Blink') omit the charset while the content is
    UTF-8, so decode the QP-decoded payload ourselves — get_content() would fall back to
    ascii-with-replacement and produce mojibake."""
    import email
    from email import policy
    msg = email.message_from_bytes(data, policy=policy.default)
    htmls, texts = [], []
    for part in msg.walk():
        ct = part.get_content_type()
        if ct not in ("text/html", "text/plain"):
            continue
        try:
            raw = part.get_payload(decode=True) or b""
            cs = part.get_content_charset() or "utf-8"
            try:
                content = raw.decode(cs, "ignore")
            except LookupError:
                content = raw.decode("utf-8", "ignore")
            (htmls if ct == "text/html" else texts).append(content)
        except Exception:  # noqa: BLE001
            continue
    if htmls:
        return html_to_text("\n".join(htmls))
    if texts:
        return "\n".join(texts).strip()[:MAX_CHARS]
    return html_to_text(data.decode("utf-8", "ignore"))


def _odf(data: bytes) -> str:
    """OpenDocument (odt/ods/odp): text of content.xml, one line per text:p element."""
    from lxml import etree
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        xml = z.read("content.xml")
    root = etree.fromstring(xml)
    ns_text = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
    lines = []
    for p in root.iter(f"{{{ns_text}}}p", f"{{{ns_text}}}h"):
        t = "".join(p.itertext()).strip()
        if t:
            lines.append(t)
        if sum(len(x) for x in lines) > MAX_CHARS:
            break
    return "\n".join(lines).strip()


def _rtf(data: bytes) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        return rtf_to_text(data.decode("utf-8", "ignore"), errors="ignore").strip()[:MAX_CHARS]
    except ImportError:
        return data.decode("utf-8", "ignore")[:MAX_CHARS]


def _default_inner(blob: bytes, iname: str) -> str:
    """Fallback member extractor when the primary one isn't supplied: plain text natively,
    everything else through our own extract()."""
    iext = iname.rsplit(".", 1)[-1].lower() if "." in iname else ""
    if iext in ("md", "markdown", "txt", "csv", "tsv", "json", "log", "yaml", "yml"):
        return blob.decode("utf-8", "ignore")
    return extract(blob, iname)


def _zip(data: bytes, inner) -> str:
    """ZIP archive: full listing + extracted text of the first readable members.
    `inner` is the primary extractor; images/binaries inside are only listed."""
    inner = inner or _default_inner
    out: list[str] = []
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        infos = [i for i in z.infolist() if not i.is_dir()]
        out.append(f"[Архив ZIP, файлов: {len(infos)}]")
        for i in infos[:60]:
            out.append(f"- {i.filename} ({i.file_size} байт)")
        if len(infos) > 60:
            out.append(f"…и ещё {len(infos) - 60}")
        done = total = 0
        for i in infos:
            if done >= _ZIP_MAX_MEMBERS or total > _ZIP_MAX_TOTAL_BYTES:
                break
            iname = i.filename
            iext = iname.rsplit(".", 1)[-1].lower() if "." in iname else ""
            readable = iext in ("pdf", "docx", "xlsx", "xlsm", "md", "markdown", "txt", "csv",
                                "tsv", "json", "log", "yaml", "yml") or iext in EXTS
            if iext == "zip" or not readable or i.file_size > _ZIP_MAX_MEMBER_BYTES:
                continue
            blob = z.read(i)
            total += len(blob)
            text = (inner(blob, iname) if inner else extract(blob, iname)) or ""
            if text.strip():
                done += 1
                budget = max(4000, (MAX_CHARS - sum(len(x) for x in out)) // max(1, _ZIP_MAX_MEMBERS - done + 1))
                out.append(f"\n## Файл из архива: {iname}\n" + text.strip()[:budget])
            if sum(len(x) for x in out) > MAX_CHARS:
                break
    return "\n".join(out).strip()
